"""
동아출판 Corporate Brand Standard (Style #31) — Sample PPTX Generator

Generates a sample presentation demonstrating all 6 slide types:
  Type A: Cover
  Type B: Table of Contents
  Type D: Section Divider
  Type C: Content (text / standard table / highlight table)
  Type E: Appendix
  E.O.D: End slide
"""

from __future__ import annotations

import os
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.table import _Cell
from pptx.util import Emu, Inches, Pt

# ── Brand Colors ──────────────────────────────────────────
BRAND_RED = RGBColor(0xD7, 0x00, 0x3F)
BRAND_DARK = RGBColor(0xC0, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)
BODY_GRAY = RGBColor(0x40, 0x40, 0x40)
MUTED_GRAY = RGBColor(0x80, 0x80, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x44, 0x72, 0xC4)
SEPARATOR_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
BORDER_GRAY = RGBColor(0xBF, 0xBF, 0xBF)
ZEBRA_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
LIGHT_RED_LINE = RGBColor(0xE8, 0xE8, 0xE8)

STATUS_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
STATUS_GREEN_TX = RGBColor(0x2E, 0x7D, 0x32)
STATUS_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE)
STATUS_BLUE_TX = RGBColor(0x15, 0x65, 0xC0)
STATUS_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
STATUS_ORANGE_TX = RGBColor(0xE6, 0x51, 0x00)
STATUS_RED_BG = RGBColor(0xFF, 0xEB, 0xEE)

# ── Font Names ────────────────────────────────────────────
NANUM_EB = "NanumGothic ExtraBold"
NANUM_REG = "NanumGothic"
MALGUN = "맑은 고딕"

# ── Slide dimensions (standard 10×7.5 inches) ────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.77)
MARGIN_R = Inches(0.77)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R  # usable content width


# ── Helpers ───────────────────────────────────────────────

def _set_font(
    run,
    name: str = MALGUN,
    size: int = 12,
    bold: bool = False,
    color: RGBColor = BLACK,
) -> None:
    """Configure run font in one call."""
    font = run.font
    font.name = name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    # East-Asian font fallback
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", name)


def _set_cell_margins(cell: _Cell, top: int = 3, bottom: int = 3, left: int = 6, right: int = 6) -> None:
    """Set cell internal margins in points."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set("marT", str(Pt(top)))
    tcPr.set("marB", str(Pt(bottom)))
    tcPr.set("marL", str(Pt(left)))
    tcPr.set("marR", str(Pt(right)))


def _set_cell_border(cell: _Cell, edges: str = "all", color: RGBColor = BORDER_GRAY, width: float = 0.75) -> None:
    """Set cell borders via XML. edges: 'all', 'outer', 'bottom', etc."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    border_names = {
        "all": ["lnL", "lnR", "lnT", "lnB"],
        "outer": ["lnL", "lnR", "lnT", "lnB"],
        "bottom": ["lnB"],
        "top": ["lnT"],
    }
    for ln_name in border_names.get(edges, ["lnL", "lnR", "lnT", "lnB"]):
        ln = tcPr.find(qn(f"a:{ln_name}"))
        if ln is not None:
            tcPr.remove(ln)
        ln = tcPr.makeelement(qn(f"a:{ln_name}"), {})
        ln.set("w", str(int(Pt(width))))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        solid = ln.makeelement(qn("a:solidFill"), {})
        srgb = solid.makeelement(qn("a:srgbClr"), {"val": str(color)})
        solid.append(srgb)
        ln.append(solid)
        tcPr.append(ln)


def _add_brand_line(slide) -> None:
    """Top-edge brand red accent line (3.5pt)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Pt(3.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_RED
    shape.line.fill.background()


def _add_separator(slide, top, width=None, color=SEPARATOR_GRAY, thickness: float = 0.5) -> None:
    """Thin horizontal separator line."""
    w = width or CONTENT_W
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, top, w, Pt(thickness),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_page_number(slide, num: int) -> None:
    """Bottom-right page number."""
    txbox = slide.shapes.add_textbox(
        Inches(8.8), Inches(6.95), Inches(0.9), Inches(0.3),
    )
    tf = txbox.text_frame
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    _set_font(run, MALGUN, 10, color=MUTED_GRAY)


def _add_slide_title(slide, title: str, top: float = 0.30) -> float:
    """Add content slide title + separator. Returns Y position below separator."""
    txbox = slide.shapes.add_textbox(
        MARGIN_L, Inches(top), CONTENT_W, Inches(0.48),
    )
    tf = txbox.text_frame
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, MALGUN, 20, bold=True, color=BLACK)

    sep_y = Inches(top) + Inches(0.52)
    _add_separator(slide, sep_y)
    return sep_y + Pt(12)


def _add_subtitle(slide, text: str, y) -> float:
    """Add a 14pt bold subtitle/subsection heading. Returns Y below."""
    txbox = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(0.35))
    tf = txbox.text_frame
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, MALGUN, 14, bold=True, color=BODY_GRAY)
    return y + Inches(0.38)


def _add_bullets(slide, bullets: Sequence[str], y, indent: float = 0.15) -> float:
    """Add bullet list starting at y. Returns Y after last bullet."""
    txbox = slide.shapes.add_textbox(
        MARGIN_L + Inches(indent), y, CONTENT_W - Inches(indent), Inches(len(bullets) * 0.38 + 0.1),
    )
    tf = txbox.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(0)
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.line_spacing = Pt(18)

        # Bullet marker
        run_dot = p.add_run()
        run_dot.text = "•   "
        _set_font(run_dot, MALGUN, 12, color=MUTED_GRAY)

        # Bullet text
        run_txt = p.add_run()
        run_txt.text = bullet
        _set_font(run_txt, MALGUN, 12, color=BODY_GRAY)

    return y + Inches(len(bullets) * 0.36 + 0.05)


def _add_key_value_line(slide, label: str, value: str, y, value_color: RGBColor = BRAND_DARK) -> float:
    """Add a single key: value line with the value in emphasis color."""
    txbox = slide.shapes.add_textbox(MARGIN_L + Inches(0.15), y, CONTENT_W - Inches(0.15), Inches(0.32))
    tf = txbox.text_frame
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]

    run_label = p.add_run()
    run_label.text = f"{label}  "
    _set_font(run_label, MALGUN, 12, bold=True, color=BODY_GRAY)

    run_val = p.add_run()
    run_val.text = value
    _set_font(run_val, MALGUN, 12, bold=True, color=value_color)

    return y + Inches(0.34)


def _asset_path(filename: str) -> str:
    return os.path.join(os.path.dirname(__file__), "..", "assets", "images", filename)


# ══════════════════════════════════════════════════════════
#  TYPE A — Cover
# ══════════════════════════════════════════════════════════
def make_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Bottom accent line ──
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), SLIDE_W, Pt(3),
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = BRAND_RED
    bottom_line.line.fill.background()

    # ── Thin separator above title area ──
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55), Inches(5.0), Pt(0.5),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = SEPARATOR_GRAY
    sep.line.fill.background()

    # ── Company logo (bottom-left) ──
    # GIF: 741×323px → aspect 2.29:1 → w=2.2" h=0.96"
    logo = _asset_path("donga_symbol.gif")
    if os.path.exists(logo):
        logo_w = Inches(1.8)
        logo_h = Inches(1.8 / 2.29)  # maintain aspect ratio ≈ 0.79"
        slide.shapes.add_picture(
            logo, Inches(0.5), Inches(6.2), logo_w, logo_h,
        )

    # ── Title ──
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(6.5), Inches(1.1))
    tf = txbox.text_frame
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(0)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.line_spacing = Pt(54)
    run = p.add_run()
    run.text = "2026년 1분기\n사업 현황 보고"
    _set_font(run, NANUM_EB, 40, bold=True, color=BLACK)

    # ── Subtitle ──
    txbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.35), Inches(6.5), Inches(0.50))
    tf2 = txbox2.text_frame
    tf2.margin_top = Pt(6)
    tf2.margin_left = Pt(0)
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "디지털 전환 및 교육 콘텐츠 혁신 전략"
    _set_font(run2, NANUM_REG, 20, color=BODY_GRAY)

    # ── Department / Date ──
    txbox3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.15), Inches(6.5), Inches(0.70))
    tf3 = txbox3.text_frame
    tf3.margin_top = Pt(6)
    tf3.margin_left = Pt(0)
    tf3.word_wrap = True

    p3a = tf3.paragraphs[0]
    p3a.space_after = Pt(3)
    run3a = p3a.add_run()
    run3a.text = "경영기획실"
    _set_font(run3a, NANUM_REG, 14, color=BODY_GRAY)

    p3b = tf3.add_paragraph()
    p3b.space_before = Pt(0)
    run3b = p3b.add_run()
    run3b.text = "2026. 03. 28"
    _set_font(run3b, NANUM_REG, 13, color=MUTED_GRAY)


# ══════════════════════════════════════════════════════════
#  TYPE B — Table of Contents
# ══════════════════════════════════════════════════════════
def make_toc(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand_line(slide)

    # Title
    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.35), CONTENT_W, Inches(0.48))
    tf = txbox.text_frame
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "목  차"
    _set_font(run, MALGUN, 20, bold=True, color=BLACK)

    _add_separator(slide, Inches(0.87))

    # Section items with colored number
    items = [
        ("01", "사업 성과 요약"),
        ("02", "디지털 전환 추진 현황"),
        ("03", "교육 콘텐츠 혁신 전략"),
        ("04", "향후 계획 및 일정"),
    ]

    y_start = Inches(1.3)
    for i, (num, label) in enumerate(items):
        y = y_start + Inches(i * 0.62)

        # Number box
        txbox_n = slide.shapes.add_textbox(Inches(1.0), y, Inches(0.6), Inches(0.40))
        tf_n = txbox_n.text_frame
        tf_n.margin_top = Pt(4)
        tf_n.margin_left = Pt(0)
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.RIGHT
        run_n = p_n.add_run()
        run_n.text = num
        _set_font(run_n, MALGUN, 16, bold=True, color=BRAND_RED)

        # Vertical dot separator
        txbox_dot = slide.shapes.add_textbox(Inches(1.68), y, Inches(0.2), Inches(0.40))
        tf_dot = txbox_dot.text_frame
        tf_dot.margin_top = Pt(3)
        p_dot = tf_dot.paragraphs[0]
        p_dot.alignment = PP_ALIGN.CENTER
        run_dot = p_dot.add_run()
        run_dot.text = "|"
        _set_font(run_dot, MALGUN, 14, color=SEPARATOR_GRAY)

        # Label
        txbox_l = slide.shapes.add_textbox(Inches(1.95), y, Inches(6.5), Inches(0.40))
        tf_l = txbox_l.text_frame
        tf_l.margin_top = Pt(5)
        tf_l.margin_left = Pt(0)
        p_l = tf_l.paragraphs[0]
        run_l = p_l.add_run()
        run_l.text = label
        _set_font(run_l, MALGUN, 15, color=BODY_GRAY)

    # Appendix separator
    appendix_y = y_start + Inches(len(items) * 0.62 + 0.25)
    _add_separator(slide, appendix_y, color=SEPARATOR_GRAY)

    txbox_a = slide.shapes.add_textbox(Inches(1.0), appendix_y + Inches(0.15), Inches(7), Inches(0.35))
    tf_a = txbox_a.text_frame
    tf_a.margin_left = Pt(0)
    p_a = tf_a.paragraphs[0]
    run_a = p_a.add_run()
    run_a.text = "[별첨]  상세 재무 데이터"
    _set_font(run_a, MALGUN, 12, color=MUTED_GRAY)

    _add_page_number(slide, 2)


# ══════════════════════════════════════════════════════════
#  TYPE D — Section Divider
# ══════════════════════════════════════════════════════════
def make_section_divider(prs: Presentation, number: int, name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section number
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.0))
    tf = txbox.text_frame
    tf.margin_left = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"0{number}" if number < 10 else str(number)
    _set_font(run, MALGUN, 52, bold=True, color=BRAND_RED)

    # Short red accent bar under number
    bar_w = Inches(0.6)
    bar_x = (SLIDE_W - bar_w) // 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_x, Inches(3.25), bar_w, Pt(3),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_RED
    shape.line.fill.background()

    # Section name
    txbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.50), Inches(9), Inches(0.55))
    tf2 = txbox2.text_frame
    tf2.margin_left = Pt(0)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = name
    _set_font(run2, MALGUN, 22, bold=True, color=BLACK)


# ══════════════════════════════════════════════════════════
#  TYPE C — Content (text + bullets)
# ══════════════════════════════════════════════════════════
def make_content_text(
    prs: Presentation,
    title: str,
    subtitle: str | None,
    bullets: Sequence[str],
    page_num: int,
    kv_lines: Sequence[tuple[str, str]] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand_line(slide)

    y = _add_slide_title(slide, title)

    if subtitle:
        y = _add_subtitle(slide, subtitle, y)

    y = _add_bullets(slide, bullets, y)

    if kv_lines:
        y += Inches(0.15)
        for label, value in kv_lines:
            y = _add_key_value_line(slide, label, value, y)

    _add_page_number(slide, page_num)


# ══════════════════════════════════════════════════════════
#  TYPE C — Content with Standard Table
# ══════════════════════════════════════════════════════════
def make_content_standard_table(prs: Presentation, page_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand_line(slide)

    y = _add_slide_title(slide, "주요 사업 지표 현황")

    # Table caption
    cap = slide.shapes.add_textbox(MARGIN_L, y, Inches(5), Inches(0.30))
    tf_cap = cap.text_frame
    tf_cap.margin_top = Pt(0)
    tf_cap.margin_left = Pt(0)
    p_cap = tf_cap.paragraphs[0]
    run_cap = p_cap.add_run()
    run_cap.text = "[표 1] 2026년 1분기 핵심 KPI"
    _set_font(run_cap, MALGUN, 11, bold=True, color=BODY_GRAY)

    y_table = y + Inches(0.35)

    headers = ["구분", "목표", "실적", "달성률", "전년 대비"]
    data = [
        ["매출액",       "150억 원",  "162억 원",   "108.0%",  "+12.3%"],
        ["신규 콘텐츠",  "50종",      "47종",       "94.0%",   "+8종"],
        ["디지털 전환율", "60.0%",     "65.2%",      "108.7%",  "+11.4%p"],
        ["고객 만족도",   "4.5점",     "4.7점",      "104.4%",  "+0.3점"],
        ["영업이익률",    "15.0%",     "17.0%",      "113.3%",  "+2.3%p"],
    ]

    rows, cols = len(data) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(
        rows, cols, MARGIN_L, y_table, CONTENT_W, Inches(rows * 0.38),
    )
    table = tbl_shape.table

    # Column widths
    col_widths = [Inches(1.8), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.86)]
    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        _set_cell_margins(cell, top=5, bottom=5, left=8, right=8)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        _set_font(run, MALGUN, 10, bold=True, color=WHITE)
        _set_cell_border(cell, "all", WHITE, 0.5)

    # Data rows
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            _set_cell_margins(cell, top=4, bottom=4, left=8, right=8)

            bg = ZEBRA_GRAY if i % 2 == 1 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

            p = cell.text_frame.paragraphs[0]
            is_first = j == 0
            p.alignment = PP_ALIGN.LEFT if is_first else PP_ALIGN.CENTER

            run = p.add_run()
            run.text = val
            _set_font(run, MALGUN, 10, bold=is_first, color=BLACK)

            inner_color = SEPARATOR_GRAY
            _set_cell_border(cell, "all", inner_color, 0.5)

    # Outer border
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            if i == 0 or i == rows - 1 or j == 0 or j == cols - 1:
                edges = []
                if i == 0:
                    edges.append("top")
                if i == rows - 1:
                    edges.append("bottom")
                for edge in edges:
                    _set_cell_border(cell, edge, BORDER_GRAY, 0.75)

    _add_page_number(slide, page_num)


# ══════════════════════════════════════════════════════════
#  TYPE C — Content with Highlight Table
# ══════════════════════════════════════════════════════════
def make_content_highlight_table(prs: Presentation, page_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand_line(slide)

    y = _add_slide_title(slide, "프로젝트 진행 현황")

    cap = slide.shapes.add_textbox(MARGIN_L, y, Inches(5), Inches(0.30))
    tf_cap = cap.text_frame
    tf_cap.margin_top = Pt(0)
    tf_cap.margin_left = Pt(0)
    p_cap = tf_cap.paragraphs[0]
    run_cap = p_cap.add_run()
    run_cap.text = "[표 2] 핵심 프로젝트 추진 상태"
    _set_font(run_cap, MALGUN, 11, bold=True, color=BODY_GRAY)

    y_table = y + Inches(0.35)

    headers = ["프로젝트명", "담당팀", "진척률", "목표일", "상태"]
    data = [
        ["AI 교과서 플랫폼",      "디지털혁신팀", "100%", "26.02.28", ("● 완료",   STATUS_GREEN_BG,  STATUS_GREEN_TX)],
        ["학습 분석 대시보드",     "데이터팀",     "72%",  "26.06.30", ("◐ 진행중", STATUS_BLUE_BG,   STATUS_BLUE_TX)],
        ["모바일 앱 리뉴얼",      "앱개발팀",     "45%",  "26.09.15", ("▲ 주의",   STATUS_ORANGE_BG, STATUS_ORANGE_TX)],
        ["ERP 시스템 전환",       "IT인프라팀",   "28%",  "26.12.31", ("■ 지연",   STATUS_RED_BG,    BRAND_DARK)],
    ]

    rows, cols = len(data) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(
        rows, cols, MARGIN_L, y_table, CONTENT_W, Inches(rows * 0.40),
    )
    table = tbl_shape.table

    col_widths = [Inches(2.2), Inches(1.5), Inches(1.2), Inches(1.5), Inches(2.06)]
    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    # Header (brand red)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_RED
        _set_cell_margins(cell, top=5, bottom=5, left=8, right=8)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        _set_font(run, MALGUN, 10, bold=True, color=WHITE)

    # Data rows
    for i, row_data in enumerate(data):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            _set_cell_margins(cell, top=4, bottom=4, left=8, right=8)
            p = cell.text_frame.paragraphs[0]

            if j == 0:
                # Project name (category column)
                cell.fill.solid()
                cell.fill.fore_color.rgb = ZEBRA_GRAY
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = row_data[j]
                _set_font(run, MALGUN, 10, bold=True, color=BLACK)

            elif j == cols - 1:
                # Status column
                status_text, bg_color, tx_color = row_data[j]
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = status_text
                _set_font(run, MALGUN, 10, bold=True, color=tx_color)

            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = row_data[j]
                _set_font(run, MALGUN, 10, color=BLACK)

            _set_cell_border(cell, "all", SEPARATOR_GRAY, 0.5)

    _add_page_number(slide, page_num)


# ══════════════════════════════════════════════════════════
#  TYPE E — Appendix
# ══════════════════════════════════════════════════════════
def make_appendix(prs: Presentation, page_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand_line(slide)

    # Appendix label
    lbl = slide.shapes.add_textbox(MARGIN_L, Inches(0.25), Inches(1.5), Inches(0.30))
    tf_lbl = lbl.text_frame
    tf_lbl.margin_top = Pt(2)
    tf_lbl.margin_left = Pt(0)
    p_lbl = tf_lbl.paragraphs[0]
    run_lbl = p_lbl.add_run()
    run_lbl.text = "[별첨]"
    _set_font(run_lbl, MALGUN, 13, color=MUTED_GRAY)

    y = _add_slide_title(slide, "상세 재무 데이터", top=0.55)

    y = _add_subtitle(slide, "매출 구성 (단위: 억 원)", y + Inches(0.05))

    bullets_revenue = [
        "교과서 사업부: 98억 원 (전년 대비 +8.9%)",
        "참고서 사업부: 42억 원 (전년 대비 +5.0%)",
        "디지털 콘텐츠: 22억 원 (전년 대비 +34.1%)",
    ]
    y = _add_bullets(slide, bullets_revenue, y)

    y = _add_subtitle(slide, "수익성 지표", y + Inches(0.15))

    bullets_profit = [
        "영업이익률: 17.0% (전년 14.7% → 2.3%p 개선)",
        "EBITDA: 31억 원 (전년 대비 +18.2%)",
        "순이익: 19억 원 (전년 대비 +15.4%)",
    ]
    y = _add_bullets(slide, bullets_profit, y)

    _add_page_number(slide, page_num)


# ══════════════════════════════════════════════════════════
#  E.O.D — End of Document
# ══════════════════════════════════════════════════════════
def make_eod(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Thin red accent line at vertical center
    bar_w = Inches(1.2)
    bar_x = (SLIDE_W - bar_w) // 2
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_x, Inches(3.15), bar_w, Pt(2),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BRAND_RED
    slide.shapes[-1].line.fill.background()

    txbox = slide.shapes.add_textbox(Inches(2.5), Inches(3.35), Inches(5), Inches(0.8))
    tf = txbox.text_frame
    tf.margin_left = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "E.O.D"
    _set_font(run, MALGUN, 24, bold=True, color=MUTED_GRAY)

    # Subtle bottom line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_x, Inches(4.05), bar_w, Pt(2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_RED
    shape.line.fill.background()


# ══════════════════════════════════════════════════════════
#  Main — Build full presentation
# ══════════════════════════════════════════════════════════
def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 표지 (Type A) ──
    make_cover(prs)

    # ── 목차 (Type B) ──
    make_toc(prs)

    # ── 섹션 1 (Type D) ──
    make_section_divider(prs, 1, "사업 성과 요약")

    # ── 콘텐츠: 핵심 성과 (Type C) ──
    make_content_text(
        prs,
        title="1분기 핵심 성과 요약",
        subtitle="전 사업부 목표 초과 달성",
        bullets=[
            "매출액 162억 원 달성 — 전년 동기 대비 12.3% 성장",
            "디지털 교과서 시장 점유율 업계 1위 유지 (38.2%)",
            "AI 기반 학습 분석 플랫폼 베타 서비스 런칭 완료 (2월)",
            "고객 만족도 4.7점 달성 — 목표 4.5점 대비 초과",
        ],
        page_num=4,
        kv_lines=[
            ("총 매출",    "162억 원 (YoY +12.3%)"),
            ("영업이익률", "17.0% (YoY +2.3%p)"),
        ],
    )

    # ── 콘텐츠: 기본형 표 (Type C) ──
    make_content_standard_table(prs, page_num=5)

    # ── 섹션 2 (Type D) ──
    make_section_divider(prs, 2, "디지털 전환 추진 현황")

    # ── 콘텐츠: 강조형 표 (Type C) ──
    make_content_highlight_table(prs, page_num=7)

    # ── 콘텐츠: 향후 계획 (Type C) ──
    make_content_text(
        prs,
        title="2분기 이후 추진 계획",
        subtitle="중점 추진 과제 4건",
        bullets=[
            "AI 학습 분석 플랫폼 정식 출시 및 학교 현장 적용 (4월)",
            "모바일 앱 UX 전면 리뉴얼 — 접근성 개선 중점 (9월)",
            "ERP 시스템 클라우드 전환 완료 및 안정화 (12월)",
            "디지털 매출 비중 30% 달성을 위한 B2C 채널 강화",
        ],
        page_num=8,
    )

    # ── 별첨 (Type E) ──
    make_appendix(prs, page_num=9)

    # ── E.O.D ──
    make_eod(prs)

    output_path = os.path.join(os.path.dirname(__file__), "style_31_corporate_brand_donga.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")


if __name__ == "__main__":
    main()
